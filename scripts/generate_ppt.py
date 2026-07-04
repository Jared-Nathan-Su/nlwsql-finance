# -*- coding: utf-8 -*-
"""NL2SQL 答辩 PPT — 18页 16:9"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

# Colors
DB=RGBColor(0x0D,0x2B,0x5E); MB=RGBColor(0x1A,0x56,0xDB); LB=RGBColor(0xE8,0xF0,0xFE)
WH=RGBColor(0xFF,0xFF,0xFF); DK=RGBColor(0x1E,0x29,0x3B); TX=RGBColor(0x33,0x40,0x55)
GR=RGBColor(0x94,0xA3,0xB8); RD=RGBColor(0xEF,0x44,0x44); GN=RGBColor(0x10,0xB9,0x81)
OG=RGBColor(0xF5,0x9E,0x0B); CB=RGBColor(0xF8,0xFA,0xFC); B2=RGBColor(0xE2,0xE8,0xF0)

def grad_bg(sl,c1,c2):
    bg=sl.background; bg.fill.gradient(); bg.fill.gradient_angle=135
    bg.fill.gradient_stops[0].color.rgb=c1; bg.fill.gradient_stops[1].color.rgb=c2
def solid_bg(sl,c):
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb=c
def box(sl,l,t,w,h,fill=None,border=None,radius=None):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    else: s.fill.background()
    if border: s.line.color.rgb=border; s.line.width=Pt(1)
    else: s.line.fill.background()
def oval(sl,l,t,w,h,color):
    o=sl.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l),Inches(t),Inches(w),Inches(h))
    o.fill.solid(); o.fill.fore_color.rgb=color; o.line.fill.background()
def txt(sl,l,t,w,h,text,size=14,color=TX,bold=False,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
    p.text=text; p.font.size=Pt(size); p.font.color.rgb=color; p.font.bold=bold; p.alignment=align
def mtxt(sl,l,t,w,h,lines,size=13,color=TX,sp=Pt(6)):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=ln; p.font.size=Pt(size); p.font.color.rgb=color; p.space_after=sp
def card(sl,l,t,w,h,title,body,tc=DB):
    box(sl,l,t,w,h,WH,B2,0.15)
    txt(sl,l+0.3,t+0.2,w-0.6,0.4,title,15,tc,True)
    mtxt(sl,l+0.3,t+0.6,w-0.6,h-0.8,body.split('\n'),11,TX,Pt(4))
def kpi(sl,l,t,v,lb,c=MB):
    box(sl,l,t,1.8,1.3,WH,B2,0.12)
    txt(sl,l,t+0.15,1.8,0.6,v,26,c,True,PP_ALIGN.CENTER)
    txt(sl,l,t+0.8,1.8,0.35,lb,9,GR,False,PP_ALIGN.CENTER)
def hdr(sl,title,sub=""):
    txt(sl,0.8,0.4,10,0.6,title,28,DB,True)
    ln=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.8),Inches(1.0),Inches(1.5),Pt(4))
    ln.fill.solid(); ln.fill.fore_color.rgb=MB; ln.line.fill.background()
    if sub: txt(sl,0.8,1.15,10,0.35,sub,12,GR)
def ftr(sl):
    box(sl,0,7.1,13.333,0.4,DB)
    txt(sl,0.5,7.12,10,0.3,"NL2SQL 财务智能问数系统 / 经营分析赛道 / 2026",9,WH)
def pn(sl,n):
    txt(sl,12.2,7.0,1,0.35,f"{n}/18",9,GR,False,PP_ALIGN.RIGHT)

# === S1 封面 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH)
box(sl,0,0,13.333,3.0,MB)
oval(sl,1.0,1.5,0.6,0.6,RGBColor(0x3A,0x82,0xF6)); oval(sl,11.5,0.5,1.0,1.0,RGBColor(0x3A,0x82,0xF6))
txt(sl,1.5,1.0,10,1.2,"NL2SQL 财务智能问数系统",44,WH,True)
txt(sl,1.5,2.2,10,0.6,"基于大模型的企业财务智能问数与经营分析平台",18,RGBColor(0xBF,0xDB,0xFE))
box(sl,1.5,3.5,2,0,border=MB); box(sl,1.5,3.55,2,0,border=MB)
mtxt(sl,1.5,3.9,6,1.2,["赛道方向：经营分析","AI 工具：DeepSeek-V4  |  LangChain  |  Streamlit","指导教师：XXX    团队：XXX"],14,TX,Pt(10))
for i,(v,l) in enumerate([("99.9%","效率提升"),("88%","SQL准确率"),("3-5s","响应速度"),("20+","问数场景")]):
    kpi(sl,1.5+i*2.8,5.6,v,l,MB)
txt(sl,5,6.9,3,0.4,"2026 年 7 月",11,GR,False,PP_ALIGN.CENTER)

# === S2 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl)
hdr(sl,"目  录","CONTENTS")
for i,(n,t,d) in enumerate([("01","业务背景与痛点分析","企业管理层的看数之痛"),
    ("02","解决方案与系统架构","NL2SQL 引擎 + 五层 Prompt 注入"),("03","核心技术解析","SQL 安全校验 / 自动纠错 / AI 分析"),
    ("04","系统演示与典型场景","3 个 Demo 展示核心能力"),("05","AI vs 传统 & 业务价值","效率提升 99%+ / ROI 700%"),
    ("06","创新点 / 不足 / 展望","四大创新 + 未来路线图")]):
    y=1.6+i*0.9; txt(sl,1.2,y,0.7,0.5,n,26,MB,True,PP_ALIGN.CENTER)
    txt(sl,2.0,y+0.02,5,0.35,t,16,DK,True); txt(sl,2.0,y+0.4,7,0.3,d,11,GR)

# === S3 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,3)
hdr(sl,"企业管理层的看数之痛","业务背景与痛点分析")
for i,s in enumerate(["管理层提问","BI分析师理解需求","多轮沟通确认","编写SQL制作报表","反馈给管理层"]):
    x=0.8+i*2.4; box(sl,x,1.8,2.0,1.0,LB,B2,0.12)
    txt(sl,x+0.1,1.95,1.8,0.7,s,11,DK,True,PP_ALIGN.CENTER)
    if i<4: txt(sl,x+2.0,2.1,0.4,0.4,">",18,MB,False,PP_ALIGN.CENTER)
for i,(v,l) in enumerate([("1-3 天","全流程耗时"),("3-5 人次","沟通成本"),("~20%","理解偏差率"),("仅 25%","员工可自助BI"),("60%","分析师花在取数")]):
    kpi(sl,0.8+i*2.4,3.1,v,l)
for i,(t,d,c) in enumerate([("响应慢","传统模式从提问到拿到数据需1-3天\n涉及多轮沟通、需求传递和手工制表\n管理层往往错过最佳决策窗口",RD),("门槛高","非技术人员必须通过BI部门才能查数据\n业务人员懂业务但不懂SQL无法自助\n数据分析师大部分时间花在取数而非分析",OG),("分析浅","传统BI只返回数据不提供业务解读\n异常指标需要人工逐个排查才能发现\n缺乏从数据到洞察到行动建议的闭环",GR)]):
    card(sl,0.8+i*4.2,4.6,3.8,2.2,t,d,c)

# === S4 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,4)
hdr(sl,"解决方案定位","从翻译官视角重新定义财务问数")
txt(sl,1.0,1.8,5,0.4,"BEFORE - 传统模式",17,RD,True)
mtxt(sl,1.0,2.3,5,3.0,["业务人员提出看数需求","BI分析师沟通确认口径和维度","分析师手动编写SQL并验证数据","出报表反馈给业务人员","前后耗时1-3天，错过决策时机","","根本矛盾: 懂业务的人不懂SQL","懂SQL的人不理解业务场景","双方反复确认 → 效率低 + 易出错"],11,TX,Pt(6))
txt(sl,7.5,1.8,5,0.4,"AFTER - AI 方案",17,GN,True)
mtxt(sl,7.5,2.3,5,3.0,["业务人员直接用自然语言输入问题","AI自动理解语义并生成SQL查询语句","系统执行查询，3-10秒返回结果","自动附带数据表格、图表和经营分析解读","异常指标自动标注，给出行动建议","","AI成为连接业务与数据的翻译官:","自然语言 → SQL（业务翻译为技术）","查询结果 → 分析洞察（数据翻译为决策）"],11,TX,Pt(6))
box(sl,1.0,5.5,11.5,0.8,LB,None,0.08)
txt(sl,1.3,5.65,11,0.5,"我们不替代财务人员，而是在业务与数据之间架一座 AI 翻译桥",14,DB,True,PP_ALIGN.CENTER)

# === S5 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,5)
hdr(sl,"系统架构","三层架构: 展示层 to NL2SQL 引擎 to 数据层")
for i,(t,lns,cl) in enumerate([("展示层 Streamlit",["自然语言输入 | 数据表格 | 可视化图表(Plotly)","SQL展示 | AI分析解读 | 历史查询记录"],MB),
    ("NL2SQL 核心引擎 LangChain+LLM",["Step1:Prompt构建 Step2:LLM生成SQL Step3:SQL校验","Step4:执行查询 Step5:AI结果解读 自动纠错(最多3次)","多模型支持: DeepSeek-V4 / Qwen / GPT-4o 一键切换"],RGBColor(0x3B,0x82,0xF6)),
    ("数据层 SQLite / 自定义上传",["6 维度表 + 5 事实表(演示数据)","支持 CSV/Excel 文件上传 自动建表自动分析","销售 | 成本 | 费用(含预算) | 应收 | 应付"],DB)]):
    y=1.8+i*1.6; box(sl,0.8,y,11.7,1.35,cl,None,0.12)
    txt(sl,1.1,y+0.08,11,0.35,t,14,WH,True)
    for j,ln in enumerate(lns): txt(sl,1.1,y+0.48+j*0.32,11,0.3,ln,10,RGBColor(0xE0,0xE8,0xF0))
for i,(v,l) in enumerate([("12 张","演示数据表"),("59,720","演示总记录"),("3 年","时间跨度"),("6 区","覆盖"),("12 线","产品线"),("任意","上传自定义")]):
    kpi(sl,0.4+i*2.15,6.65,v,l)

# === S6 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,6)
hdr(sl,"数据如何组织？—— 数据库设计","围绕企业经营的 6 个维度 + 5 类业务数据")
for i,(n,cn,d) in enumerate([("时间维度","按年/季/月/日分析","覆盖2022-2024三年数据"),("产品维度","12条产品线","硬件、软件、服务三大类"),("部门维度","8个部门","研发、销售、市场、财务等"),("区域维度","全国6城市4大区","华东、华南、华北、西南"),("客户维度","50个客户","含信用评级和行业标签"),("供应商","20家供应商","覆盖主要采购渠道")]):
    y=1.7+i*0.48; box(sl,0.8,y,5.5,0.4,LB if i%2==0 else WH,B2,0.04)
    txt(sl,0.9,y+0.04,1.8,0.32,n,9,DB,True); txt(sl,2.8,y+0.04,3.5,0.32,f"{cn} / {d}",9,TX)
for i,(n,cn,r,d) in enumerate([("销售记录","收入+销量","20,736条","每笔销售的时间/产品/区域"),("成本记录","直接+间接成本","18,412条","与销售收入精细匹配"),("费用记录","预算+实际","17,280条","各部门预算执行情况追踪"),("应收记录","客户欠款","1,644条","含是否逾期和逾期天数"),("应付记录","公司欠款","446条","供应商付款到期日管理")]):
    y=1.7+i*0.48; box(sl,6.8,y,5.7,0.4,LB if i%2==0 else WH,B2,0.04)
    txt(sl,6.9,y+0.04,1.8,0.32,n,9,DB,True); txt(sl,8.8,y+0.04,1.2,0.32,r,8,GR); txt(sl,10.1,y+0.04,2.2,0.32,f"{cn}:{d}",9,TX)
txt(sl,0.8,4.9,11,0.4,"系统能回答哪些财务问题？",14,DB,True)
mtxt(sl,0.8,5.3,11.5,2.0,["营收: 某个区域、产品线、时间段赚了多少钱？跟去年同期比涨了还是跌了？",
    "盈利: 哪个产品最赚钱？哪个区域的毛利率在下降？需要重点关注什么地方？",
    "花钱: 各部门花了多少钱？有没有超预算？哪类费用增长最快？",
    "回款: 哪些客户欠款还没还？欠了多少？超过约定期限多久了？该催谁？",
    "综合: 上个月公司整体经营情况怎么样？收入、利润、费用、回款都正常吗？"],10,TX,Pt(5))

# === S7 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,7)
hdr(sl,"AI 怎样准确理解财务问题？","不是简单把问题丢给AI，而是分五层逐步引导，准确率从72%提升到88%")
for i,(nm,title,desc,reason) in enumerate([("第1步","告诉AI它是谁","你是一位资深财务数据分析师\n精通企业经营指标和SQL","让AI进入角色"),
    ("第2步","告诉AI有哪些数据","数据库有12张表，每张表存什么信息\n字段名称和含义是什么","防止AI瞎编字段"),("第3步","告诉AI怎么算","毛利率是(收入-成本)/收入\n同比是和去年同期比较","确保计算逻辑正确"),
    ("第4步","给AI看几个例子","类似问题应该生成什么样的SQL\n给3-5个标准问答示范","让AI模仿正确格式"),("第5步","提出用户真正的问题","华东区第二季度毛利率\n和去年同期相比变化多少？","结合前面知识给出答案")]):
    y=1.6+i*0.85; cl=[MB,RGBColor(0x3B,0x82,0xF6),RGBColor(0x63,0x66,0xF1),OG,GN][i]
    box(sl,0.8,y,0.55,0.65,cl,None,0.08); txt(sl,0.8,y+0.12,0.55,0.4,str(i+1),16,WH,True,PP_ALIGN.CENTER)
    txt(sl,1.5,y+0.02,2,0.3,title,12,cl,True); txt(sl,1.5,y+0.3,8,0.5,desc,10,TX); txt(sl,9.8,y+0.1,3,0.3,f"目的: {reason}",9,GR)
box(sl,0.8,6.0,11.7,0.8,LB,None,0.08)
mtxt(sl,1.0,6.1,11,0.6,["效果验证: 直接让AI生成SQL准确率仅72% | 五层逐步引导后达到88% | 再加更多示例可达91%"],11,DB,Pt(4))

# === S8 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,8)
hdr(sl,"如何保证AI生成的SQL安全可靠？","AI可能犯错，我们设了4道防线 + 最多3次自动修正")
for i,s in enumerate(["语法检查\nSQL写法对不对","权限检查\n只能查不能改","表名检查\n数据表存在吗","执行验证\n跑一遍看结果"]):
    x=0.8+i*3.1; box(sl,x,1.7,2.7,1.0,LB,MB,0.12)
    txt(sl,x+0.1,1.85,2.5,0.7,s,12,DB,True,PP_ALIGN.CENTER)
    if i<3: txt(sl,x+2.7,2.0,0.4,0.4,">",18,GN,False,PP_ALIGN.CENTER)
for i,(t,d) in enumerate([("操作权限控制","仅开放SELECT查询权限，拦截DROP/DELETE/UPDATE等危险操作"),
    ("表名字段校验","白名单机制验证SQL中的表名和字段名，防止AI幻觉产生不存在的对象"),("自动纠错闭环","校验失败时将错误信息反馈给LLM，引导其修正SQL并重新生成，最多重试3次"),
    ("实战效果","超过80%的错误在3次重试内自动修复，用户无需关心底层纠错过程")]):
    y=3.2+i*0.75; box(sl,0.8,y,11.7,0.6,CB,None,0.06)
    txt(sl,1.0,y+0.1,4,0.4,t,12,DB,True); txt(sl,5.0,y+0.1,7,0.4,d,10,TX)

# === S9 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,9)
hdr(sl,"覆盖 20+ 典型财务问数场景","从简单查询到复杂多维分析 + 自定义数据即插即用")
for i,(t,d) in enumerate([("营收查询","2024年总营收?\n各区域排名?季度趋势?"),("毛利分析","各产品线毛利率排名\n华东区Q2同比变化?"),
    ("费用管控","预算执行率<80%部门\n费用占比最大类型"),("客户分析","营收Top5客户\n逾期应收最多客户"),
    ("风险预警","逾期超30天应收\n毛利率下降产品线"),("自定义数据","上传CSV/Excel即用\n自动建表+智能问数")]):
    x=0.5+(i%3)*4.2; y=1.6+(i//3)*2.3; card(sl,x,y,3.8,2.0,t,d)

# === S10 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,10)
hdr(sl,"系统演示 - 4 个典型场景","实际运行效果 | SQL自动生成+AI解读 | 平均响应3-5秒")
for i,(t,lns) in enumerate([("场景1:毛利率同比分析",["输入:2024年Q2毛利率同比2023年Q2变化?","SQL:自动生成含CTE跨年对比查询","结果:2024Q2to32.50%|2023Q2to34.60%|down2.10%","AI解读:主因原材料成本up8.3%华南区价格竞争","响应:3.2秒"]),
    ("场景2:预算执行监控",["输入:哪些部门预算执行率低于80%?","SQL:HAVING子句+三表JOIN聚合","结果:研发中心62.3%|市场部71.5%","AI解读:建议确认研发里程碑评估预算调整","响应:2.8秒"]),
    ("场景3:应收风险预警",["输入:逾期超30天应收款按客户统计","SQL:多条件过滤+聚合+排序","结果:12个客户逾期总计2,847,000元","AI解读:客户_03逾期680K超90天to立即催收","响应:3.5秒"]),
    ("场景4:自定义数据上传",["上传:拖入CSV/Excel员工表/销售表","自动:建表+Schema生成+快捷问题","提问:员工平均薪资?各部门人数?","AI解读:基于任意上传数据智能分析","响应:4.0秒"])]):
    card(sl,0.2+i*3.25,1.7,3.0,3.2,t,'\n'.join(lns))
for i,(v,l) in enumerate([("3-5秒","平均响应"),("88%","SQL准确率"),("80%+","纠错成功率"),("自定义","数据即插即用")]):
    kpi(sl,1.0+i*3.1,5.5,v,l,GN)

# === S11 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,11)
hdr(sl,"AI vs 传统方案 - 为什么非AI不可?","全方位对比 + 不可替代性论证")
hd=["对比维度","传统BI方案","AI方案","提升"]; ws=[3.5,3.0,3.0,2.5]
for j,h in enumerate(hd):
    x=0.8+sum(ws[:j]); box(sl,x,1.7,ws[j],0.5,DB); txt(sl,x+0.1,1.78,ws[j]-0.2,0.35,h,12,WH,True,PP_ALIGN.CENTER)
for i,rw in enumerate([["使用方式","写SQL/拖拽配置","自然语言输入","零门槛"],["响应时间","1-3天","3-10秒","99.9%up"],
    ["问题覆盖率","仅预设模板","开放域/泛化强","无限"],["数据解读","只给数据","数据+分析+建议","质的飞跃"],
    ["异常发现","人工抽查","AI自动标注","不遗漏"],["部署成本","BI系统+人力","API+开源框架","极低"]]):
    for j,cell in enumerate(rw):
        x=0.8+sum(ws[:j]); bg=CB if i%2==0 else WH; box(sl,x,2.2+i*0.48,ws[j],0.48,bg,B2)
        co=GN if j==3 else (MB if j==2 else TX); txt(sl,x+0.1,2.3+i*0.48,ws[j]-0.2,0.32,cell,10,co,False,PP_ALIGN.CENTER)
box(sl,0.8,5.3,11.7,1.6,LB,None,0.08)
mtxt(sl,1.0,5.4,11,1.4,["为什么非AI不可?",
    "1.自然语言的多样性:同一问题有几十种问法->传统规则无法穷举只有LLM能理解语义",
    "2.复杂查询的组合性:用户可能问出从未预定义的多维交叉问题->LLM可动态理解并生成SQL",
    "3.分析解读的专业性:不只给数据像财务分析师一样给出解读->只有LLM具备此能力"],10,DB,Pt(5))

# === S12 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,12)
hdr(sl,"AI 在方案中的三重角色","不只是工具，而是贯穿整个问数流程的智能引擎")
for i,(t,d,c) in enumerate([("智能翻译官","将自然语言精准转换为SQL查询\n理解业务术语如毛利率、同比、预算执行率\n自动识别时间范围、区域、产品线等维度\n华东区毛利率 → SELECT...JOIN...WHERE",MB),
    ("数据分析师","查询结果自动生成经营分析解读\n包含: 数据摘要、变动原因分析\n影响因素拆解和风险提示\n毛利率下降2.1%主因原材料成本上升8.3%",GN),
    ("质量守门员","四道防线确保生成的SQL安全可执行\n语法校验 → 权限检查 → 字段验证 → 执行测试\n校验失败自动反馈LLM修正，最多重试3次\n80%+的错误在重试中被自动修复",OG)]):
    x=0.8+i*4.2; box(sl,x,1.7,3.8,3.2,WH,c,0.15)
    txt(sl,x+0.3,1.85,3.2,0.5,t,18,c,True); mtxt(sl,x+0.4,2.5,3.0,2.2,d.split('\n'),12,TX,Pt(10))
txt(sl,3.5,5.4,6,0.5,"AI 让财务人员从取数工升级为分析决策者",15,DB,True,PP_ALIGN.CENTER)

# === S13 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,13)
hdr(sl,"业务价值 - 定量分析","可量化的效率成本与准确率提升")
for i,(v,l) in enumerate([("99.9%","问数响应提速"),("88%","SQL生成准确率"),("70%+","BI重复取数减少"),("3-5s","平均响应时间"),("210K","年节省人力成本"),("80%+","纠错成功率")]):
    kpi(sl,0.5+i*2.1,1.7,v,l)
box(sl,0.8,3.3,11.7,2.8,CB,B2,0.12); txt(sl,1.1,3.4,11,0.4,"投入产出测算（保守估计）",15,DB,True)
mtxt(sl,1.1,3.85,11,2.1,["测算前提: 10人管理团队，每人每月平均查询5次经营数据","",
    "传统模式成本: 每次查询平均耗时2小时（沟通+写SQL+出报表）× 人力成本200元/小时 ≈ 20,000元/月 → 240,000元/年",
    "AI方案成本: 大模型API调用约500元/月 + 系统运维约2,000元/月 ≈ 2,500元/月 → 30,000元/年","",
    "年度净节省: 240,000 - 30,000 = 210,000元    投资回报率(ROI): 700%",
    "注: 以上仅为可量化的人力成本，不含决策提速带来的间接业务价值"],11,TX,Pt(4))

# === S14 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,14)
hdr(sl,"业务价值 - 定性维度","超越数字的战略意义，改变企业数据使用方式")
for i,(t,d,c) in enumerate([("决策敏捷性","从会前准备报表到会中实时问数\n管理层可以当场追问、当场验证假设\n决策节奏从月/周级缩短到分钟级\n典型场景: 经营分析会中直接问毛利率变化原因",GN),
    ("数据民主化","每个部门经理都可以用自然语言查数据\n不再依赖BI团队排期，消除数据瓶颈\n让数据从少数人的工具变为全员可用\n释放组织各级人员的数据分析潜能",MB),
    ("风险防控前置","从月底出报表的事后总结变为实时监控\nAI自动识别异常指标并主动预警\n逾期应收、预算超支、毛利下滑即时发现\n风险发现从事后变为事中，甚至事前预警",RD),
    ("知识数字化","AI每次分析的解读都是可复用的知识资产\n新人无需长时间学习即可获得专业数据分析能力\n降低对资深财务分析师个人经验的依赖\n企业数据分析和财务知识持续积累沉淀",OG)]):
    y=1.8+i*1.05; box(sl,0.8,y,11.7,0.9,WH,c,0.12)
    txt(sl,1.1,y+0.08,3,0.5,t,16,c,True); mtxt(sl,1.1,y+0.55,10.5,0.6,d.split('\n'),12,TX)

# === S15 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,15)
hdr(sl,"四大创新点","技术创新驱动业务价值")
for i,(t,d) in enumerate([("五层 Prompt 注入策略","业界首创角色toSchemato规则to示例to问题分层注入准确率提升15个百分点"),
    ("SQL 自动纠错闭环","语法校验to权限检查to字段验证to执行验证失败自动反馈LLM修正成功率80%+"),
    ("问数+分析一体化","LLM二次调用自动生成经营分析解读:数据摘要+变动原因+风险提示+行动建议"),
    ("财务场景深度适配","内置毛利率/净利率/同比环比/预算执行率等专业指标覆盖20+问数模式支持多模型热替换"),
    ("自定义数据即插即用","上传CSV/Excel自动建表动态Schema生成表名白名单注入一键分析任意数据")]):
    y=1.7+i*1.25; box(sl,0.8,y,11.7,1.05,LB if i%2==0 else CB,None,0.08)
    txt(sl,1.1,y+0.1,10,0.32,t,15,DB,True); mtxt(sl,1.1,y+0.48,10.5,0.55,d.split('\n'),11,TX)

# === S16 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,16)
hdr(sl,"不足与展望","正视当前局限，规划未来演进路线")
txt(sl,1.0,1.7,5.5,0.5,"当前不足",16,RD,True)
mtxt(sl,1.0,2.3,5.5,4.0,["复杂多步推理准确率偏低: 涉及跨表嵌套查询时，","  AI有时难以正确分解步骤，需进一步优化Prompt策略","模糊问题澄清能力有限: 用户输入过于简单时，","  系统无法像人类分析师一样主动追问确认口径","模拟数据局限性: 当前基于模拟企业数据验证，","  缺乏真实企业业务场景的多样性和复杂性测试","自定义数据Schema理解: 对用户上传的任意格式","  数据，字段语义识别和关联推断能力有待加强"],10,TX,Pt(5))
txt(sl,7.0,1.7,5.5,0.5,"未来展望",16,GN,True)
mtxt(sl,7.0,2.3,5.5,4.0,["Agent模式: 支持多轮追问、自动拆解复杂问题，","  像人类分析师一样逐步深入挖掘数据","预测分析: 接入时间序列预测模型，","  从描述过去升级为预测未来趋势","知识库RAG: 接入企业财务制度和行业知识文档，","  让AI分析更贴合企业实际业务场景","多模态输出: 支持自动生成PPT简报和图文报告，","  将分析结果直接转换为管理层可用的汇报材料","企业级部署: 对接ERP系统实现实时数据同步，","  支持多用户权限管理和数据访问控制","私有化部署: 支持完全离线部署，","  敏感财务数据不出企业内网"],10,TX,Pt(5))

# === S17 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH); ftr(sl); pn(sl,17)
hdr(sl,"答辩三问 - 标准回答","每组必须清晰回答的三个核心问题")
for i,(q,a) in enumerate([("1.解决什么财务业务问题?痛点在哪里?",
    "管理层看数难看数慢-传统问数需1-3天依赖BI团队。三大痛点:1响应慢(天级);2门槛高(需SQL技能);3分析浅(只给数据不解读)。核心矛盾:业务人员懂业务但不懂SQL技术人员懂SQL但不懂业务。"),
    ("2.AI在你的方案中扮演什么角色?为什么非AI不可?",
    "三重角色:智能翻译官(自然语言toSQL)+数据分析师(结果to解读)+质量守门员(校验纠错)。非AI不可:1自然语言有无限多样性传统规则无法穷举;2复杂多维交叉查询无法预定义;3只有大模型能结合财务知识给出专业分析解读。"),
    ("3.你的方案能带来什么业务价值?",
    "效率:问数从天级to秒级提升99.9%。准确率:SQL生成88%避免人工翻译错误。风控:AI实时监控to异常不遗漏风险发现从事后变事中。ROI:年节省人力成本约210,000元投资回报率700%。定性:决策敏捷性+数据民主化+知识数字化沉淀。")]):
    y=1.7+i*1.75; box(sl,0.8,y,11.7,1.55,CB,MB,0.08)
    txt(sl,1.1,y+0.1,11,0.35,q,13,DB,True); mtxt(sl,1.1,y+0.5,11,1.0,a.split('\n'),10,TX,Pt(3))

# === S18 致谢 ===
sl=prs.slides.add_slide(prs.slide_layouts[6]); solid_bg(sl,WH)
box(sl,0,4.5,13.333,3.0,MB)
oval(sl,1.5,5.5,1.0,1.0,RGBColor(0x3A,0x82,0xF6)); oval(sl,11.5,5.0,0.7,0.7,RGBColor(0x3A,0x82,0xF6))
txt(sl,2,1.0,9,1.0,"感谢聆听",48,DB,True,PP_ALIGN.CENTER)
txt(sl,2,2.0,9,0.6,"NL2SQL 财务智能问数系统",20,GR,False,PP_ALIGN.CENTER)
box(sl,5,2.7,3,0,border=MB)
txt(sl,2,3.0,9,0.5,"让数据会说话，让决策有依据",15,TX,False,PP_ALIGN.CENTER)
mtxt(sl,2,5.2,9,1.5,["GitHub: github.com/Jared-Nathan-Su/nlwsql-finance","Streamlit Cloud: nlwsql-finance.streamlit.app","多模型: DeepSeek-V4 / Qwen / GPT-4o","支持: 演示数据 + CSV/Excel自定义上传","欢迎各位评委老师提问！"],13,WH,Pt(8))
for p in sl.shapes[-1].text_frame.paragraphs: p.alignment=PP_ALIGN.CENTER

# ===== SAVE =====
out=os.path.join(os.path.dirname(os.path.dirname(__file__)),"docs","NL2SQL_答辩PPT.pptx")
prs.save(out)
print(f"OK: {out} ({len(prs.slides)} slides)")
